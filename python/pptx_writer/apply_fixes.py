"""Apply ACME auto-fixes to a PPTX while preserving run/cell formatting."""

from __future__ import annotations

import re
from io import BytesIO
from typing import Any, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

TITLE_FONT_SIZE_PT = 24
BULLET_FONT_SIZE_PT = 16
TABLE_HEADER_COLOR = "006EBE"

ACME_PALETTE = {
    "006EBE",
    "0070C0",
    "FFFF00",
    "FFFFDB",
    "FFFFFF",
    "000000",
    "595959",
    "7F7F7F",
    "FF0000",
}


def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.replace("#", "").upper()
    if len(h) == 8:
        h = h[2:]
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _normalize_hex(hex_color: Optional[str]) -> Optional[str]:
    if not hex_color:
        return None
    h = hex_color.replace("#", "").upper()
    if len(h) == 8:
        h = h[2:]
    return h if len(h) == 6 else None


def _color_distance(a: str, b: str) -> float:
    ar, ag, ab = int(a[0:2], 16), int(a[2:4], 16), int(a[4:6], 16)
    br, bg, bb = int(b[0:2], 16), int(b[2:4], 16), int(b[4:6], 16)
    return ((ar - br) ** 2 + (ag - bg) ** 2 + (ab - bb) ** 2) ** 0.5


def _nearest_palette(hex_color: str) -> str:
    n = _normalize_hex(hex_color) or hex_color
    best = n
    best_dist = float("inf")
    for palette in ACME_PALETTE:
        d = _color_distance(n, palette)
        if d < best_dist:
            best_dist = d
            best = palette
    return best


def _set_run_color(run, hex_color: str) -> None:
    run.font.color.rgb = _hex_to_rgb(hex_color)


def _set_cell_fill(cell, hex_color: str) -> None:
    fill = cell.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(hex_color)


def _iter_text_frames(slide):
    if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
        yield slide.shapes.title.text_frame
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            yield shape.text_frame
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        yield cell.text_frame


def _replace_in_text_frame(
    text_frame, pattern: str | re.Pattern[str], replacement: str
) -> bool:
    changed = False
    regex = (
        re.compile(re.escape(pattern), re.IGNORECASE)
        if isinstance(pattern, str)
        else pattern
    )
    for para in text_frame.paragraphs:
        for run in para.runs:
            before = run.text
            if regex.search(before):
                run.text = regex.sub(replacement, before)
                if run.text != before:
                    changed = True
    return changed


def _replace_across_slide(
    slide, pattern: str | re.Pattern[str], replacement: str
) -> bool:
    changed = False
    for text_frame in _iter_text_frames(slide):
        changed = _replace_in_text_frame(text_frame, pattern, replacement) or changed
    return changed


def _apply_title_font(
    slide,
    *,
    family: Optional[str] = None,
    bold: Optional[bool] = None,
    size_pt: Optional[float] = None,
) -> bool:
    title = slide.shapes.title
    if title is None or not title.has_text_frame:
        return False
    touched = False
    for para in title.text_frame.paragraphs:
        for run in para.runs:
            if not run.text.strip():
                continue
            if family is not None:
                run.font.name = family
            if bold is not None:
                run.font.bold = bold
            if size_pt is not None:
                run.font.size = Pt(size_pt)
            touched = True
    return touched


def _apply_footer_font(slide, *, family: str, size_pt: float) -> bool:
    touched = False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.lower()
        if not any(k in text for k in ("confidential", "proprietary", "privacy")):
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.name = family
                run.font.size = Pt(size_pt)
                touched = True
    return touched


def _apply_bullet_font_size(slide, size_pt: float) -> bool:
    touched = False
    for shape in slide.shapes:
        if not shape.has_text_frame or shape == slide.shapes.title:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    run.font.size = Pt(size_pt)
                    touched = True
    return touched


def _apply_bullet_font_family(slide, family: str) -> bool:
    touched = False
    for shape in slide.shapes:
        if not shape.has_text_frame or shape == slide.shapes.title:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    run.font.name = family
                    touched = True
    return touched


def _apply_table_header_fill(slide, hex_color: str) -> bool:
    touched = False
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        table = shape.table
        header_rows = min(2, len(table.rows))
        for r_idx in range(header_rows):
            for cell in table.rows[r_idx].cells:
                _set_cell_fill(cell, hex_color)
                touched = True
    return touched


def _apply_color_snap(slide) -> bool:
    touched = False
    for text_frame in _iter_text_frames(slide):
        for para in text_frame.paragraphs:
            for run in para.runs:
                try:
                    if run.font.color.type is not None and run.font.color.rgb is not None:
                        rgb = run.font.color.rgb
                        hex_val = f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}".upper()
                        if hex_val not in ACME_PALETTE:
                            _set_run_color(run, _nearest_palette(hex_val))
                            touched = True
                except Exception:
                    pass
    for shape in slide.shapes:
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    try:
                        if cell.fill.type is not None and cell.fill.fore_color.rgb is not None:
                            rgb = cell.fill.fore_color.rgb
                            hex_val = f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}".upper()
                            if hex_val not in ACME_PALETTE:
                                _set_cell_fill(cell, _nearest_palette(hex_val))
                                touched = True
                    except Exception:
                        pass
    return touched


def apply_finding_to_slide(slide, finding: dict[str, Any]) -> bool:
    rule_id = finding.get("rule_id", "")
    actual = (finding.get("actual_value") or "").strip()
    expected = (finding.get("expected_value") or "").strip()

    if rule_id == "TITLE_001":
        return _apply_title_font(slide, family="Calibri", bold=True)
    if rule_id == "TITLE_002":
        return _apply_title_font(slide, size_pt=TITLE_FONT_SIZE_PT)
    if rule_id == "TITLE_003":
        title = slide.shapes.title
        if title is None or not title.has_text_frame:
            return False
        changed = False
        for para in title.text_frame.paragraphs:
            for run in para.runs:
                before = run.text
                run.text = re.sub(r"[.!?]+\s*$", "", before)
                if run.text != before:
                    changed = True
        return changed
    if rule_id == "TITLE_005":
        return _replace_across_slide(slide, re.compile(r"\bDRAFT\b", re.I), "")
    if rule_id == "BULLET_001":
        return _apply_bullet_font_size(slide, BULLET_FONT_SIZE_PT)
    if rule_id == "BULLET_002":
        changed = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    before = run.text
                    run.text = re.sub(r"[.!?:;]+\s*$", "", before)
                    if run.text != before:
                        changed = True
        return changed
    if rule_id == "BULLET_004":
        return _apply_bullet_font_family(slide, "Calibri")
    if rule_id == "COLOR_001":
        return _apply_color_snap(slide)
    if rule_id == "TABLE_001":
        return _apply_table_header_fill(slide, TABLE_HEADER_COLOR)
    if rule_id == "TABLE_003":
        return _replace_across_slide(
            slide, re.compile(r"\bpercentile\b", re.I), "%ile"
        )
    if rule_id == "TERM_001":
        return _replace_across_slide(slide, re.compile(r"\bTGT\b"), "Target")
    if rule_id == "TERM_002":
        return _replace_across_slide(
            slide, re.compile(r"Company\s+Name", re.I), "Company"
        )
    if rule_id == "TERM_003":
        changed = False
        for variant in [v.strip() for v in actual.split(";") if v.strip()]:
            short = variant.split()[0] if variant.split() else ""
            if short and short != variant:
                changed = _replace_across_slide(slide, variant, short) or changed
        return changed
    if rule_id == "TERM_004":
        wrong_m = re.search(r'"([^"]+)"', actual)
        target_m = re.search(r'"([^"]+)"', expected)
        if wrong_m and target_m:
            wrong = wrong_m.group(1)
            target = target_m.group(1)
            return _replace_across_slide(
                slide, re.compile(rf"\b{re.escape(wrong)}\b", re.I), target
            )
        return False
    if rule_id == "FOOTER_002":
        return _apply_footer_font(slide, family="Calibri", size_pt=8)
    if rule_id == "FOOTER_004":
        touched = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.lower()
            if not any(k in text for k in ("confidential", "proprietary", "privacy")):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.bold = False
                    run.font.italic = False
                    touched = True
        return touched

    if (
        actual
        and expected
        and len(actual) < 80
        and len(expected) < 80
        and actual != expected
        and "ACME palette" not in expected
        and not expected.startswith("≤")
    ):
        return _replace_across_slide(slide, actual.split(";")[0].strip(), expected)

    return False


def apply_fixes_to_pptx(
    file_bytes: bytes, findings: list[dict[str, Any]]
) -> tuple[bytes, int, int]:
    prs = Presentation(BytesIO(file_bytes))
    applied = 0
    skipped = 0

    for finding in findings:
        slide_number = int(finding.get("slide_number", 0))
        if slide_number < 1 or slide_number > len(prs.slides):
            skipped += 1
            continue
        slide = prs.slides[slide_number - 1]
        try:
            if apply_finding_to_slide(slide, finding):
                applied += 1
            else:
                skipped += 1
        except Exception:
            skipped += 1

    out = BytesIO()
    prs.save(out)
    return out.getvalue(), applied, skipped
