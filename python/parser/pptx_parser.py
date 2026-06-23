"""Extract normalized metadata from PPTX files. No validation logic."""

from __future__ import annotations

import re
from io import BytesIO
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Pt

from python.schemas.metadata import (
    ChartMetadata,
    ColorMetadata,
    ParagraphMetadata,
    PositionMetadata,
    PresentationMetadata,
    ShapeMetadata,
    SlideMetadata,
    SlideType,
    TableCellMetadata,
    TableMetadata,
    TextMetadata,
    TextRunMetadata,
)


EMU_PER_INCH = 914400


def _emu_to_inches(emu: int) -> float:
    return round(emu / EMU_PER_INCH, 4)


def _position(shape) -> PositionMetadata:
    return PositionMetadata(
        left_inches=_emu_to_inches(shape.left),
        top_inches=_emu_to_inches(shape.top),
        width_inches=_emu_to_inches(shape.width),
        height_inches=_emu_to_inches(shape.height),
    )


def _rgb_to_hex(rgb: RGBColor) -> str:
    return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def _extract_color(color_format) -> Optional[ColorMetadata]:
    try:
        if color_format is None:
            return None
        if color_format.type is not None and hasattr(color_format, "rgb"):
            rgb = color_format.rgb
            if rgb is not None:
                hex_val = _rgb_to_hex(rgb)
                return ColorMetadata(hex=hex_val, rgb=str(rgb))
    except Exception:
        return None
    return None


def _extract_fill_hex(shape) -> Optional[str]:
    try:
        fill = shape.fill
        if fill.type is not None and hasattr(fill, "fore_color"):
            color = _extract_color(fill.fore_color)
            if color and color.hex:
                return color.hex.upper()
    except Exception:
        pass
    return None


def _bullet_char(paragraph) -> Optional[str]:
    try:
        p_pr = paragraph._p.pPr  # noqa: SLF001
        if p_pr is None:
            return None
        bu_char = p_pr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}buChar")
        if bu_char is not None:
            return bu_char.get("char")
        bu_auto = p_pr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum")
        if bu_auto is not None:
            return "auto"
    except Exception:
        pass
    return None


def _paragraph_indent_inches(paragraph) -> Optional[float]:
    """Explicit left margin (marL) of the paragraph, or None when inherited.

    We deliberately return None when marL is not set on the paragraph itself,
    rather than guessing from the level. Indentation rules only validate
    explicit values so inherited master indentation never produces a false
    positive.
    """
    try:
        pf = paragraph._p.pPr  # noqa: SLF001
        if pf is None:
            return None
        mar_l = pf.get("marL")
        if mar_l is not None:
            return round(int(mar_l) / EMU_PER_INCH, 4)
    except Exception:
        pass
    return None


def _paragraph_line_spacing(paragraph) -> Optional[float]:
    """Line spacing as a multiple (e.g. 1.0 = single).

    python-pptx returns a float for "multiple" spacing and a Length for exact
    point spacing. We only capture the multiple form, since the guideline
    ("single line spacing") is expressed as a multiple. Exact-point spacing is
    left as None to avoid mis-comparing units.
    """
    try:
        raw = paragraph.line_spacing
        if isinstance(raw, (int, float)):
            return round(float(raw), 3)
    except Exception:
        pass
    return None


DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _optional_bool(value) -> Optional[bool]:
    if value is None:
        return None
    return bool(value)


def _resolve_def_rpr(shape, level: int = 0) -> dict:
    """Read inherited default run properties from shape lstStyle."""
    result: dict = {}
    try:
        tx_body = shape._element.txBody  # noqa: SLF001
        if tx_body is None:
            return result
        lst_style = tx_body.find(f"{{{DRAWING_NS}}}lstStyle")
        if lst_style is None:
            return result

        lvl_pr = lst_style.find(f"{{{DRAWING_NS}}}lvl{level + 1}pPr")
        if lvl_pr is None:
            lvl_pr = lst_style.find(f"{{{DRAWING_NS}}}defPPr")

        def_rpr = None
        if lvl_pr is not None:
            def_rpr = lvl_pr.find(f"{{{DRAWING_NS}}}defRPr")
        if def_rpr is None:
            def_rpr = lst_style.find(f"{{{DRAWING_NS}}}defPPr")

        if def_rpr is None:
            return result

        b_attr = def_rpr.get("b")
        if b_attr is not None:
            result["bold"] = b_attr in ("1", "true")

        sz = def_rpr.get("sz")
        if sz is not None:
            result["size_pt"] = round(int(sz) / 100, 2)

        latin = def_rpr.find(f"{{{DRAWING_NS}}}latin")
        if latin is not None and latin.get("typeface"):
            result["family"] = latin.get("typeface")
    except Exception:
        pass
    return result


def _run_superscript(run) -> Optional[bool]:
    """True when the run's baseline is positive (rendered as superscript)."""
    try:
        rpr = run._r.rPr  # noqa: SLF001
        if rpr is None:
            return None
        baseline = rpr.get("baseline")
        if baseline is None:
            return None
        return int(baseline) > 0
    except Exception:
        return None


def _extract_run_metadata(run, shape, paragraph_level: int) -> TextRunMetadata:
    font = run.font
    family = font.name
    size_pt = round(font.size.pt, 2) if font.size is not None else None
    bold = _optional_bool(font.bold)
    italic = _optional_bool(font.italic)

    if family is None or bold is None or size_pt is None:
        defaults = _resolve_def_rpr(shape, paragraph_level)
        if family is None:
            family = defaults.get("family")
        if size_pt is None:
            size_pt = defaults.get("size_pt")
        if bold is None and "bold" in defaults:
            bold = defaults["bold"]

    return TextRunMetadata(
        text=run.text,
        font_family=family,
        font_size_pt=size_pt,
        bold=bold,
        italic=italic,
        color=_extract_color(font.color),
        superscript=_run_superscript(run),
    )


def _paragraph_alignment(paragraph) -> Optional[str]:
    """Explicit horizontal alignment name (LEFT/CENTER/RIGHT/JUSTIFY)."""
    try:
        algn = paragraph.alignment
        if algn is None:
            return None
        name = getattr(algn, "name", None)
        return name
    except Exception:
        return None


def _vertical_anchor(text_frame) -> Optional[str]:
    """Explicit vertical anchor name (TOP/MIDDLE/BOTTOM) of a text frame/cell."""
    try:
        anchor = text_frame.vertical_anchor
        if anchor is None:
            return None
        return getattr(anchor, "name", None)
    except Exception:
        return None


def _extract_paragraphs(text_frame, shape) -> list[ParagraphMetadata]:
    paragraphs: list[ParagraphMetadata] = []
    for para in text_frame.paragraphs:
        level = para.level or 0
        runs = [_extract_run_metadata(run, shape, level) for run in para.runs]
        space_after = None
        try:
            if para.space_after is not None:
                space_after = round(para.space_after.pt, 2)
        except Exception:
            pass
        paragraphs.append(
            ParagraphMetadata(
                level=para.level or 0,
                text=para.text,
                runs=runs,
                bullet_char=_bullet_char(para),
                indent_inches=_paragraph_indent_inches(para),
                space_after_pt=space_after,
                line_spacing=_paragraph_line_spacing(para),
                alignment=_paragraph_alignment(para),
            )
        )
    return paragraphs


def _placeholder_info(shape) -> tuple[bool, Optional[str]]:
    """Return (is_placeholder, placeholder_type_name) for a shape."""
    try:
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            name = getattr(ph_type, "name", None) or str(ph_type)
            return True, name
    except Exception:
        pass
    return False, None


def _extract_text_shape(shape, is_title: bool = False) -> TextMetadata:
    tf = shape.text_frame
    paragraphs = _extract_paragraphs(tf, shape)
    full_text = shape.text.strip()
    is_placeholder, placeholder_type = _placeholder_info(shape)
    return TextMetadata(
        shape_id=str(shape.shape_id),
        shape_name=shape.name,
        is_title=is_title,
        is_placeholder=is_placeholder,
        placeholder_type=placeholder_type,
        position=_position(shape),
        paragraphs=paragraphs,
        full_text=full_text,
        vertical_anchor=_vertical_anchor(tf),
        fill_hex=_extract_fill_hex(shape),
    )


def _cell_fill_hex(cell) -> Optional[str]:
    try:
        fill = cell.fill
        if fill.type is not None:
            color = _extract_color(fill.fore_color)
            if color and color.hex:
                return color.hex.upper()
    except Exception:
        pass
    return None


def _cell_font_color(cell) -> Optional[str]:
    try:
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                color = _extract_color(run.font.color)
                if color and color.hex:
                    return color.hex.upper()
    except Exception:
        pass
    return None


def _is_merged_cell(cell) -> bool:
    try:
        tc = cell._tc  # noqa: SLF001
        grid_span = tc.get("gridSpan")
        row_span = tc.get("rowSpan")
        v_merge = tc.find("{http://schemas.openxmlformats.org/drawingml/2006/main}vMerge")
        h_merge = tc.find("{http://schemas.openxmlformats.org/drawingml/2006/main}hMerge")
        return bool(grid_span or row_span or v_merge is not None or h_merge is not None)
    except Exception:
        return False


_DML_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _cell_alignment(cell) -> Optional[str]:
    try:
        para = cell.text_frame.paragraphs[0]
        return _paragraph_alignment(para)
    except Exception:
        return None


def _cell_italic(cell) -> Optional[bool]:
    try:
        para = cell.text_frame.paragraphs[0]
        if para.runs:
            return _optional_bool(para.runs[0].font.italic)
    except Exception:
        pass
    return None


def _cell_border(cell, edge: str) -> Optional[bool]:
    """Whether the cell has an explicit visible border on the given edge.

    ``edge`` is one of ``lnL``/``lnR``/``lnT``/``lnB``. Returns True when a line
    element with a real fill is present, False when an explicit "no line" is set,
    and None when the edge is inherited (unknown).
    """
    try:
        tc = cell._tc  # noqa: SLF001
        tc_pr = tc.find(f"{_DML_NS}tcPr")
        if tc_pr is None:
            return None
        ln = tc_pr.find(f"{_DML_NS}{edge}")
        if ln is None:
            return None
        # Explicit "no fill" line means the border is intentionally absent.
        if ln.find(f"{_DML_NS}noFill") is not None:
            return False
        # A solidFill (or any fill child) indicates a visible border.
        if ln.find(f"{_DML_NS}solidFill") is not None:
            return True
        # Width without noFill is treated as a visible border.
        if ln.get("w"):
            return True
        return None
    except Exception:
        return None


def _extract_table(shape) -> TableMetadata:
    table = shape.table
    cells: list[TableCellMetadata] = []
    header_fill = None

    for r_idx, row in enumerate(table.rows):
        for c_idx, cell in enumerate(row.cells):
            font_size = None
            font_family = None
            bold: Optional[bool] = None
            try:
                para = cell.text_frame.paragraphs[0]
                if para.runs:
                    run = para.runs[0]
                    font_family = run.font.name
                    if run.font.size:
                        font_size = round(run.font.size.pt, 2)
                    bold = _optional_bool(run.font.bold)
            except Exception:
                pass
            fill = _cell_fill_hex(cell)
            if r_idx == 0 and fill:
                header_fill = fill
            cells.append(
                TableCellMetadata(
                    row=r_idx,
                    col=c_idx,
                    text=cell.text.strip(),
                    font_family=font_family,
                    font_size_pt=font_size,
                    bold=bold,
                    italic=_cell_italic(cell),
                    fill_hex=fill,
                    font_color_hex=_cell_font_color(cell),
                    is_merged=_is_merged_cell(cell),
                    alignment=_cell_alignment(cell),
                    vertical_anchor=_vertical_anchor(cell.text_frame),
                    border_left=_cell_border(cell, "lnL"),
                    border_right=_cell_border(cell, "lnR"),
                    border_top=_cell_border(cell, "lnT"),
                    border_bottom=_cell_border(cell, "lnB"),
                )
            )

    return TableMetadata(
        shape_id=str(shape.shape_id),
        shape_name=shape.name,
        position=_position(shape),
        rows=len(table.rows),
        cols=len(table.columns),
        cells=cells,
        header_row_fill_hex=header_fill,
    )


def _extract_chart_or_picture(shape) -> ChartMetadata:
    is_picture = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    meta = ChartMetadata(
        shape_id=str(shape.shape_id),
        shape_name=shape.name,
        position=_position(shape),
        is_picture=is_picture,
    )
    if hasattr(shape, "chart") and not is_picture:
        try:
            chart = shape.chart
            meta.is_native_chart = True
            if chart.has_title and chart.chart_title:
                meta.has_title = True
                meta.title_text = chart.chart_title.text_frame.text.strip()
                try:
                    tf = chart.chart_title.text_frame
                    runs = []
                    for para in tf.paragraphs:
                        runs.extend(para.runs)
                    sizes = [r.font.size.pt for r in runs if r.font.size]
                    if sizes:
                        meta.title_font_size_pt = round(max(sizes), 2)
                    bold_vals = [_optional_bool(r.font.bold) for r in runs if r.text.strip()]
                    explicit = [b for b in bold_vals if b is not None]
                    if explicit:
                        meta.title_bold = any(explicit)
                except Exception:
                    pass
        except Exception:
            pass
        try:
            meta.has_legend = bool(shape.chart.has_legend)
        except Exception:
            meta.has_legend = None
        try:
            axes = []
            for axis_attr in ("category_axis", "value_axis"):
                try:
                    axis = getattr(shape.chart, axis_attr)
                    axes.append(bool(getattr(axis, "has_title", False)))
                except Exception:
                    continue
            meta.has_axis_titles = any(axes) if axes else None
        except Exception:
            meta.has_axis_titles = None
    return meta


def _classify_slide(slide, slide_number: int) -> SlideType:
    texts = [sh.text.strip().lower() for sh in slide.shapes if hasattr(sh, "text") and sh.text.strip()]
    if slide_number == 1:
        return SlideType.TITLE
    if len(slide.shapes) <= 3 and any(len(t) < 60 for t in texts) and not any(
        hasattr(sh, "table") for sh in slide.shapes
    ):
        for sh in slide.shapes:
            if hasattr(sh, "text") and sh.text.strip() and not hasattr(sh, "table"):
                if sh.top < Emu(2000000):
                    return SlideType.SECTION
    return SlideType.CONTENT


def _find_confidentiality(texts: list[TextMetadata]) -> tuple[Optional[TextMetadata], bool]:
    keywords = ("confidential", "proprietary", "privacy")
    for text in texts:
        lower = text.full_text.lower()
        if any(k in lower for k in keywords):
            return text, True
    for text in texts:
        if text.position.top_inches > 6.5 and len(text.full_text) < 200:
            if text.paragraphs:
                for para in text.paragraphs:
                    for run in para.runs:
                        if run.font_size_pt and run.font_size_pt <= 9:
                            return text, True
    return None, False


def _find_title(texts: list[TextMetadata], slide) -> Optional[TextMetadata]:
    if slide.shapes.title is not None and slide.shapes.title.text.strip():
        return _extract_text_shape(slide.shapes.title, is_title=True)
    candidates = sorted(
        [t for t in texts if t.position.top_inches < 2.0 and len(t.full_text) > 0],
        key=lambda t: t.position.top_inches,
    )
    return candidates[0] if candidates else None


def _parse_authors_and_date(slide, slide_type: SlideType) -> tuple[list[str], Optional[str], bool]:
    authors: list[str] = []
    date_text = None
    contains_draft = False
    if slide_type != SlideType.TITLE:
        return authors, date_text, contains_draft
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        text = shape.text.strip()
        if not text:
            continue
        if re.search(r"\bDRAFT\b", text, re.IGNORECASE):
            contains_draft = True
        date_match = re.search(
            r"\b(?:January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{1,2},?\s+\d{4}\b",
            text,
            re.IGNORECASE,
        )
        if date_match:
            date_text = date_match.group(0)
        if re.search(r"\bauthor", text, re.IGNORECASE):
            authors.append(text)
        elif re.search(r"\bprepared\s+by\b", text, re.IGNORECASE):
            authors.append(text)
        elif (
            shape.top > Emu(1500000)
            and len(text) < 120
            and not date_match
            and not re.search(r"\bDRAFT\b", text, re.IGNORECASE)
            and re.search(r"[A-Z][a-z]+\s+[A-Z][a-z]+", text)
            and not re.search(r"\b(confidential|proprietary|acme)\b", text, re.IGNORECASE)
        ):
            authors.append(text)
    return authors, date_text, contains_draft


def parse_pptx(file_bytes: bytes, filename: str) -> PresentationMetadata:
    prs = Presentation(BytesIO(file_bytes))
    slides_meta: list[SlideMetadata] = []

    for idx, slide in enumerate(prs.slides, start=1):
        slide_type = _classify_slide(slide, idx)
        texts: list[TextMetadata] = []
        tables: list[TableMetadata] = []
        charts: list[ChartMetadata] = []
        shapes: list[ShapeMetadata] = []

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                tables.append(_extract_table(shape))
            elif shape.shape_type in (MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.PICTURE):
                charts.append(_extract_chart_or_picture(shape))
            elif hasattr(shape, "text_frame"):
                if shape == slide.shapes.title:
                    continue
                texts.append(_extract_text_shape(shape))

        title = _find_title(texts, slide)
        if title and title in texts:
            texts = [t for t in texts if t.shape_id != title.shape_id]

        all_texts = ([title] if title else []) + texts
        confidentiality, has_conf = _find_confidentiality(all_texts)
        if confidentiality and confidentiality in all_texts:
            texts = [t for t in texts if t.shape_id != confidentiality.shape_id]

        notes = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass

        authors, date_text, contains_draft = _parse_authors_and_date(slide, slide_type)

        for shape in slide.shapes:
            if shape.shape_type not in (
                MSO_SHAPE_TYPE.TABLE,
                MSO_SHAPE_TYPE.CHART,
                MSO_SHAPE_TYPE.PICTURE,
            ) and not (hasattr(shape, "text_frame") and shape.text.strip()):
                font_size = None
                font_family = None
                bold = None
                shape_text = None
                try:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        shape_text = shape.text.strip() or None
                        if shape.text_frame.paragraphs:
                            para = shape.text_frame.paragraphs[0]
                            if para.runs:
                                run = para.runs[0]
                                font_family = run.font.name
                                if run.font.size:
                                    font_size = round(run.font.size.pt, 2)
                                bold = _optional_bool(run.font.bold)
                except Exception:
                    pass
                shapes.append(
                    ShapeMetadata(
                        shape_id=str(shape.shape_id),
                        shape_name=shape.name,
                        shape_type=str(shape.shape_type),
                        position=_position(shape),
                        fill_hex=_extract_fill_hex(shape),
                        text=shape_text,
                        font_family=font_family,
                        font_size_pt=font_size,
                        bold=bold,
                    )
                )

        slides_meta.append(
            SlideMetadata(
                slide_number=idx,
                slide_type=slide_type,
                notes=notes,
                title=title,
                confidentiality=confidentiality,
                texts=texts,
                tables=tables,
                charts=charts,
                shapes=shapes,
                has_confidentiality=has_conf,
                authors=authors,
                date_text=date_text,
                contains_draft=contains_draft,
            )
        )

    return PresentationMetadata(
        filename=filename,
        slide_width_inches=_emu_to_inches(prs.slide_width),
        slide_height_inches=_emu_to_inches(prs.slide_height),
        slide_count=len(slides_meta),
        slides=slides_meta,
    )
